"""
File Processor for ClarityOS
Extracts text from uploaded files: PDF, DOCX, CSV, XLSX, PPTX, TXT, MD
"""
import io
from fastapi import UploadFile

async def extract_text_from_file(file: UploadFile) -> str:
    """
    Extract text content from various file types.
    Returns the text as a string.
    """
    filename = file.filename.lower()
    content = await file.read()
    
    try:
        # PDF files
        if filename.endswith(".pdf"):
            import pypdf
            pdf = pypdf.PdfReader(io.BytesIO(content))
            text = ""
            for page in pdf.pages:
                text += page.extract_text() + "\n"
            return text
        
        # Word documents    
        elif filename.endswith(".docx") or filename.endswith(".doc"):
            import docx
            doc = docx.Document(io.BytesIO(content))
            return "\n".join([para.text for para in doc.paragraphs])
        
        # Plain text and Markdown
        elif filename.endswith(".txt") or filename.endswith(".md"):
            return content.decode("utf-8")
        
        # CSV files
        elif filename.endswith(".csv"):
            import pandas as pd
            df = pd.read_csv(io.BytesIO(content))
            return df.to_string()
        
        # Excel files
        elif filename.endswith(".xlsx") or filename.endswith(".xls"):
            import pandas as pd
            df = pd.read_excel(io.BytesIO(content))
            return df.to_string()
        
        # PowerPoint files
        elif filename.endswith(".ppt") or filename.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            text = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text.append(shape.text)
            return "\n".join(text)
        
        else:
            return f"Unsupported file format: {filename}. Supported: PDF, DOCX, TXT, MD, CSV, XLSX, PPTX"
            
    except Exception as e:
        print(f"Error extracting text from {filename}: {e}")
        return f"Error reading file: {str(e)}"
